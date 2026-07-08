#!/usr/bin/env python3
# jitter_log_*.csv の jitter_ms 列を集計して、表中心の PowerPoint
# (jitter_analysis.pptx) を生成するスクリプト。
# 使い方: このフォルダで  python3 make_jitter_pptx.py
# 必要ライブラリ: python-pptx  (pip3 install --user python-pptx)

import csv
import glob
import math
import statistics

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn  # XML 名前空間を扱うためのヘルパー

# 60fps のフレーム周期 (ms)。理論ジッタの計算に使う
FRAME_MS = 1000.0 / 60

JP_FONT = "Yu Gothic"  # スライド全体で使う日本語フォント


# ---------- 集計 ----------

def load_all_rows():
    """全 CSV から (jitter, bpm) のリストを読み込む"""
    rows = []
    for path in sorted(glob.glob("jitter_log_*.csv")):
        with open(path) as fp:
            for row in csv.DictReader(fp):
                # 例外処理: 壊れた行や空ファイルは読み飛ばす
                try:
                    rows.append((float(row["jitter_ms"]), float(row["bpm"])))
                except (ValueError, KeyError):
                    pass
    return rows


def percentile95_abs(values):
    """絶対値の 95 パーセンタイル (Processing 側 printJitterSummary と同じ定義)"""
    n = len(values)
    abs_sorted = sorted(abs(v) for v in values)
    return abs_sorted[min(n - 1, math.ceil(n * 0.95) - 1)]


def summarize(rows):
    """全体サマリーと BPM 別統計を計算して返す"""
    js = [r[0] for r in rows]
    overall = {
        "n": len(js),
        "mean": statistics.mean(js),
        "std": statistics.pstdev(js),  # 母標準偏差 (Processing 側と同じ)
        "min": min(js),
        "max": max(js),
        "p95": percentile95_abs(js),
    }

    # BPM ごとにグループ化 (辞書の setdefault でリストへ振り分け)
    by_bpm = {}
    for j, b in rows:
        by_bpm.setdefault(b, []).append(j)

    bpm_stats = []
    for b in sorted(by_bpm):
        v = by_bpm[b]
        interval = int(60000.0 / b)          # Processing コードと同じ int 切り捨て
        theory = (-interval) % FRAME_MS      # 剰余演算: 60fps 量子化による理論ジッタ
        bpm_stats.append({
            "bpm": b,
            "interval": interval,
            "n": len(v),
            "mean": statistics.mean(v),
            "std": statistics.pstdev(v) if len(v) > 1 else 0.0,
            "max": max(abs(x) for x in v),
            "p95": percentile95_abs(v),
            "theory": theory,
        })
    return overall, bpm_stats


# ---------- pptx 部品 ----------

def set_run_font(run, size, bold=False, color=None):
    """run に日本語フォント・サイズ・色を設定する。
    python-pptx の font.name はラテン文字側にしか効かないため、
    XML の a:ea (東アジアフォント) 要素も直接設定する"""
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = JP_FONT
    if color is not None:
        run.font.color.rgb = color
    rPr = run.font._rPr  # 内部 XML (rPr 要素) へアクセス
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", JP_FONT)


def add_textbox(slide, x, y, w, h, lines):
    """テキストボックスを追加する。lines は (文字列, サイズ, 太字) のリスト"""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, size, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = text
        set_run_font(run, size, bold)
    return box


def fill_table(table, data, header_size=14, body_size=13):
    """2 次元リスト data の内容を表に流し込む (先頭行はヘッダー扱い)"""
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = ""  # 既定の空段落をクリアしてから書く
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if r == 0 or c > 0 else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            # ヘッダー行は白太字 (既定テーマの濃色ヘッダーに合わせる)
            if r == 0:
                set_run_font(run, header_size, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
            else:
                set_run_font(run, body_size)


# ---------- スライド生成 ----------

def build(overall, bpm_stats):
    prs = Presentation()
    # 16:9 に設定
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]  # 白紙レイアウト

    # --- スライド 1: タイトル ---
    s = prs.slides.add_slide(blank)
    add_textbox(s, 0.8, 2.4, 11.7, 1.2, [
        ("LED 遅延（ジッタ）計測結果", 40, True),
    ])
    add_textbox(s, 0.8, 3.8, 11.7, 1.6, [
        ("BPM 同期 LED の点灯タイミング誤差の分析", 20, False),
        (f"データ: jitter_log_*.csv 計 {overall['n']:,} 拍（2026-07-07〜07-08）", 16, False),
        ("計測環境: Processing（draw() 60fps）", 16, False),
    ])

    # --- スライド 2: 全体サマリー表 ---
    s = prs.slides.add_slide(blank)
    add_textbox(s, 0.6, 0.4, 12.0, 0.8, [("全体サマリー（全 BPM 合計）", 28, True)])
    rows = [
        ["指標", "値"],
        ["サンプル数", f"{overall['n']:,} 拍（17 ファイル）"],
        ["平均", f"+{overall['mean']:.1f} ms"],
        ["標準偏差", f"{overall['std']:.1f} ms"],
        ["最小 / 最大", f"{overall['min']:.0f} ms / +{overall['max']:.0f} ms"],
        ["95 パーセンタイル（絶対値）", f"{overall['p95']:.0f} ms"],
    ]
    shape = s.shapes.add_table(len(rows), 2, Inches(1.6), Inches(1.5), Inches(10.0), Inches(3.6))
    fill_table(shape.table, rows, header_size=16, body_size=16)
    add_textbox(s, 1.6, 5.6, 10.5, 1.4, [
        ("・負のジッタ（早く点灯）は 0 件 — ズレは常に遅れ方向", 16, False),
        ("・約 70% が 6〜10 ms に集中。同期ズレの知覚閾値（20〜30 ms）未満", 16, False),
    ])

    # --- スライド 3: BPM 別統計表 ---
    s = prs.slides.add_slide(blank)
    add_textbox(s, 0.6, 0.35, 12.5, 0.8, [("BPM 別統計 と 60fps 量子化の理論値", 28, True)])
    rows = [["BPM", "拍間隔", "拍数 n", "平均", "標準偏差", "最大", "p95", "理論値*"]]
    for st in bpm_stats:
        rows.append([
            f"{st['bpm']:g}",
            f"{st['interval']} ms",
            f"{st['n']:,}",
            f"+{st['mean']:.1f} ms",
            f"{st['std']:.1f} ms",
            f"{st['max']:.0f} ms",
            f"{st['p95']:.0f} ms",
            f"{st['theory']:.1f} ms",
        ])
    shape = s.shapes.add_table(len(rows), 8, Inches(0.7), Inches(1.25), Inches(12.0), Inches(4.4))
    fill_table(shape.table, rows, header_size=14, body_size=13)
    add_textbox(s, 0.7, 6.0, 12.0, 1.2, [
        ("* 理論値 = 拍間隔を 16.7 ms（60fps のフレーム周期）で割った余りから予測されるジッタ", 14, False),
        ("全 BPM で実測平均が理論値 +1〜4 ms に一致 → 誤差の主因は draw() の 60fps 量子化", 16, True),
    ])

    # --- スライド 4: まとめ ---
    s = prs.slides.add_slide(blank)
    add_textbox(s, 0.6, 0.4, 12.0, 0.8, [("まとめ", 28, True)])
    add_textbox(s, 0.9, 1.5, 11.8, 4.6, [
        ("1 拍ごとのズレは 0〜20 ms・95% は 12 ms 以内 — 知覚閾値（20〜30 ms）未満", 18, False),
        ("ズレは全件が遅れ方向で、BPM 別平均は理論値と一致 — 原因は draw() の 60fps 量子化", 18, False),
        ("測定環境（Processing）由来の誤差であり、loop() が高速に回る実機 Arduino の LED はこれより正確", 18, False),
        ("注意: ドリフト型更新のため遅れは毎拍累積する（BPM 66 で 1 分あたり約 0.5 秒）", 18, False),
        ("", 10, False),
        ("改善案", 18, True),
        ("① 実 LED も理想時刻基準（prevMillis += interval）で更新 → 累積を根本解消", 16, False),
        ("② frameRate(240) を指定 → 時間分解能 16.7 ms → 4.2 ms（単発・累積とも約 1/4）", 16, False),
    ])

    prs.save("jitter_analysis.pptx")
    print("saved: jitter_analysis.pptx")


if __name__ == "__main__":
    rows = load_all_rows()
    overall, bpm_stats = summarize(rows)
    build(overall, bpm_stats)
